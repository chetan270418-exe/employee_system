"""
EmpTrack Cloud Computing Internship PPT Generator
Generates a professional, cloud-themed PowerPoint presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# ── Color Palette (Cloud/AWS Theme) ──────────────────────────────────
DARK_BG       = RGBColor(0x0F, 0x1B, 0x2D)   # Deep navy
MEDIUM_BG     = RGBColor(0x15, 0x23, 0x3A)   # Card bg
LIGHT_BG      = RGBColor(0x1A, 0x2C, 0x47)   # Slightly lighter
ACCENT_ORANGE = RGBColor(0xFF, 0x99, 0x00)   # AWS orange
ACCENT_BLUE   = RGBColor(0x23, 0x8B, 0xE6)   # Cloud blue
ACCENT_TEAL   = RGBColor(0x00, 0xD4, 0xAA)   # Teal green
ACCENT_PURPLE = RGBColor(0x7C, 0x3A, 0xED)   # Purple accent
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xBB, 0xC7, 0xD9)
MID_GRAY      = RGBColor(0x8A, 0x9B, 0xB0)
RED_ACCENT    = RGBColor(0xEF, 0x44, 0x44)
GREEN_ACCENT  = RGBColor(0x22, 0xC5, 0x5E)

# ── Slide width/height ───────────────────────────────────────────────
SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    """Set a solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """Add a shape to a slide."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.background()  # transparent by default
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    if line_width is not None:
        shape.line.width = Pt(line_width)
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18, font_color=WHITE,
                bold=False, alignment=PP_ALIGN.LEFT, font_name='Segoe UI'):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide_text(text_frame, items, font_size=16, font_color=LIGHT_GRAY, font_name='Segoe UI', spacing=Pt(8)):
    """Add bullet point items to an existing text frame."""
    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = font_name
        p.space_after = spacing
        p.level = 0


def add_accent_line(slide, left, top, width, color=ACCENT_ORANGE, height=Pt(4)):
    """Add a thin colored accent line."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, fill_color=MEDIUM_BG, border_color=None):
    """Add a rounded rectangle card."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_decorative_circles(slide):
    """Add subtle decorative circles to the slide background."""
    # Top-right circle
    c1 = add_shape(slide, MSO_SHAPE.OVAL, Inches(11), Inches(-0.5), Inches(3), Inches(3))
    c1.fill.solid()
    c1.fill.fore_color.rgb = RGBColor(0x23, 0x8B, 0xE6)
    c1.fill.fore_color.brightness = 0.0
    # Make it semi-transparent via a very dark shade
    c1.fill.fore_color.rgb = RGBColor(0x14, 0x28, 0x45)

    # Bottom-left circle
    c2 = add_shape(slide, MSO_SHAPE.OVAL, Inches(-1), Inches(5.5), Inches(3), Inches(3))
    c2.fill.solid()
    c2.fill.fore_color.rgb = RGBColor(0x18, 0x2A, 0x42)


# ═════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═════════════════════════════════════════════════════════════════════

def slide_title(prs):
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, DARK_BG)
    add_decorative_circles(slide)

    # Top label
    add_textbox(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.5),
                'CLOUD COMPUTING INTERNSHIP PROJECT', font_size=14, font_color=ACCENT_ORANGE,
                bold=True, alignment=PP_ALIGN.CENTER)

    # Main title
    add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
                'Cloud-Based Smart Employee Attendance\n& Payroll Management System',
                font_size=36, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Accent line
    add_accent_line(slide, Inches(5), Inches(3.3), Inches(3.3), ACCENT_ORANGE)

    # Subtitle
    add_textbox(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.8),
                'Deployed on AWS EC2 with RDS MySQL & S3 Bucket Integration',
                font_size=20, font_color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Tech badges row
    techs = ['Python Flask', 'AWS EC2', 'AWS RDS', 'AWS S3', 'MySQL', 'Nginx']
    badge_width = Inches(1.6)
    start_x = Inches(1.5)
    for i, tech in enumerate(techs):
        x = start_x + i * (badge_width + Inches(0.15))
        card = add_card(slide, x, Inches(5.0), badge_width, Inches(0.55), LIGHT_BG, ACCENT_BLUE)
        add_textbox(slide, x, Inches(5.05), badge_width, Inches(0.5),
                    tech, font_size=12, font_color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

    # Presented by
    add_textbox(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.4),
                'Presented by: Chetan  |  Internship: Cloud Computing',
                font_size=14, font_color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(6.6), Inches(11), Inches(0.4),
                'August 2026',
                font_size=12, font_color=MID_GRAY, alignment=PP_ALIGN.CENTER)


def slide_agenda(prs):
    """Slide 2: Agenda / Table of Contents"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_decorative_circles(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
                'AGENDA', font_size=14, font_color=ACCENT_ORANGE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_ORANGE)
    add_textbox(slide, Inches(0.8), Inches(1.15), Inches(8), Inches(0.6),
                'What We Will Cover', font_size=30, font_color=WHITE, bold=True)

    agenda_items = [
        ('01', 'Project Overview & Objectives', 'What is EmpTrack and why we built it'),
        ('02', 'System Architecture & Tech Stack', 'Flask, MySQL, AWS infrastructure design'),
        ('03', 'AWS Services Used', 'EC2, RDS, S3 — What and Why'),
        ('04', 'Phase 1: Local Development', 'Building the application on localhost'),
        ('05', 'Phase 2: AWS EC2 Setup', 'Launching & configuring the cloud server'),
        ('06', 'Phase 3: AWS RDS Setup', 'Creating the managed MySQL database'),
        ('07', 'Phase 4: AWS S3 Setup', 'Setting up file/document storage bucket'),
        ('08', 'Phase 5: Deployment & Go-Live', 'Nginx, Gunicorn, Security Groups'),
        ('09', 'Security & Best Practices', 'How we secured the system on AWS'),
        ('10', 'Live Demo & Results', 'Working application walkthrough'),
    ]

    col1_items = agenda_items[:5]
    col2_items = agenda_items[5:]

    for col_idx, items in enumerate([col1_items, col2_items]):
        x_base = Inches(0.8) + col_idx * Inches(6.2)
        for i, (num, title, desc) in enumerate(items):
            y = Inches(2.2) + i * Inches(1.0)
            # Number badge
            badge = add_card(slide, x_base, y, Inches(0.55), Inches(0.55), ACCENT_ORANGE)
            add_textbox(slide, x_base, y + Inches(0.05), Inches(0.55), Inches(0.45),
                        num, font_size=16, font_color=DARK_BG, bold=True, alignment=PP_ALIGN.CENTER)
            # Title
            add_textbox(slide, x_base + Inches(0.7), y, Inches(4.5), Inches(0.35),
                        title, font_size=16, font_color=WHITE, bold=True)
            # Description
            add_textbox(slide, x_base + Inches(0.7), y + Inches(0.35), Inches(4.5), Inches(0.3),
                        desc, font_size=11, font_color=MID_GRAY)


def slide_project_overview(prs):
    """Slide 3: Project Overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_decorative_circles(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
                'PROJECT OVERVIEW', font_size=14, font_color=ACCENT_ORANGE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_ORANGE)
    add_textbox(slide, Inches(0.8), Inches(1.15), Inches(10), Inches(0.6),
                'What is EmpTrack?', font_size=30, font_color=WHITE, bold=True)

    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(0.8),
                'EmpTrack is a cloud-based Smart Employee Attendance & Payroll Management System designed\n'
                'to automate HR operations including attendance tracking, leave management, payroll processing,\n'
                'and document management — all deployed on AWS cloud infrastructure.',
                font_size=15, font_color=LIGHT_GRAY)

    # Key features cards (3 columns)
    features = [
        ('Smart Attendance', 'QR Code scanning, GPS\nlocation verification, and\nmanual override options', ACCENT_BLUE),
        ('Automated Payroll', 'Indian tax regime calculations\nPF, ESI, TDS deductions\nPDF salary slip generation', ACCENT_ORANGE),
        ('Role-Based Access', 'Admin, HR, and Employee\ndashboards with specific\npermissions and workflows', ACCENT_TEAL),
    ]
    for i, (title, desc, color) in enumerate(features):
        x = Inches(0.8) + i * Inches(4.1)
        card = add_card(slide, x, Inches(3.4), Inches(3.8), Inches(2.2), MEDIUM_BG, color)
        add_textbox(slide, x + Inches(0.3), Inches(3.6), Inches(3.2), Inches(0.4),
                    title, font_size=18, font_color=color, bold=True)
        add_accent_line(slide, x + Inches(0.3), Inches(4.05), Inches(1.5), color, Pt(2))
        add_textbox(slide, x + Inches(0.3), Inches(4.2), Inches(3.2), Inches(1.2),
                    desc, font_size=13, font_color=LIGHT_GRAY)

    # Question box
    q_card = add_card(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(1.0), LIGHT_BG, ACCENT_PURPLE)
    add_textbox(slide, Inches(1.2), Inches(6.1), Inches(10), Inches(0.35),
                'Why Cloud-Based?', font_size=16, font_color=ACCENT_PURPLE, bold=True)
    add_textbox(slide, Inches(1.2), Inches(6.45), Inches(10), Inches(0.4),
                'Traditional on-premise HR systems are costly, hard to scale, and require manual maintenance. '
                'Cloud deployment on AWS provides scalability, reliability, pay-as-you-go pricing, and global accessibility.',
                font_size=12, font_color=LIGHT_GRAY)


def slide_objectives(prs):
    """Slide 4: Project Objectives"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_decorative_circles(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
                'OBJECTIVES', font_size=14, font_color=ACCENT_ORANGE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_ORANGE)
    add_textbox(slide, Inches(0.8), Inches(1.15), Inches(10), Inches(0.6),
                'What We Aimed to Achieve', font_size=30, font_color=WHITE, bold=True)

    objectives = [
        ('Automate Attendance Tracking', 'Replace manual attendance registers with smart QR code scanning and GPS-based location verification to eliminate proxy attendance.'),
        ('Automated Payroll Processing', 'Calculate monthly salaries automatically considering working days, overtime, leaves, PF (12%), ESI (0.75%), and Indian income tax slabs.'),
        ('Cloud Deployment on AWS', 'Deploy the application on AWS EC2, connect to a managed RDS MySQL database, and use S3 for secure document storage.'),
        ('Role-Based Access Control', 'Implement Admin, HR, and Employee roles with specific dashboards and permission-gated features.'),
        ('Scalable & Secure Architecture', 'Design the system using AWS VPC security groups, Nginx reverse proxy, and Gunicorn WSGI server for production-grade deployment.'),
    ]

    for i, (title, desc) in enumerate(objectives):
        y = Inches(2.0) + i * Inches(1.05)
        # Number circle
        badge = add_shape(slide, MSO_SHAPE.OVAL, Inches(0.8), y, Inches(0.45), Inches(0.45), ACCENT_ORANGE)
        add_textbox(slide, Inches(0.8), y + Inches(0.03), Inches(0.45), Inches(0.4),
                    str(i+1), font_size=14, font_color=DARK_BG, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.5), y, Inches(3.5), Inches(0.35),
                    title, font_size=16, font_color=WHITE, bold=True)
        add_textbox(slide, Inches(1.5), y + Inches(0.35), Inches(10.5), Inches(0.6),
                    desc, font_size=12, font_color=LIGHT_GRAY)


def slide_architecture(prs):
    """Slide 5: System Architecture Diagram"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_decorative_circles(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
                'SYSTEM ARCHITECTURE', font_size=14, font_color=ACCENT_ORANGE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_ORANGE)
    add_textbox(slide, Inches(0.8), Inches(1.15), Inches(10), Inches(0.6),
                'How the System is Designed', font_size=30, font_color=WHITE, bold=True)

    # User box
    user_card = add_card(slide, Inches(0.5), Inches(3.0), Inches(2.2), Inches(1.5), LIGHT_BG, ACCENT_TEAL)
    add_textbox(slide, Inches(0.5), Inches(3.15), Inches(2.2), Inches(0.35),
                'End Users', font_size=16, font_color=ACCENT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.55), Inches(2.2), Inches(0.8),
                'Admin / HR / Employee\nWeb Browser Access\n(Desktop & Mobile)', font_size=11, font_color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Arrow 1
    arrow1 = add_shape(slide, MSO_SHAPE.RIGHT_ARROW, Inches(2.8), Inches(3.5), Inches(0.9), Inches(0.4), ACCENT_ORANGE)

    # AWS Cloud box (large)
    aws_box = add_card(slide, Inches(3.8), Inches(2.0), Inches(9.0), Inches(5.0), RGBColor(0x12, 0x20, 0x38), ACCENT_ORANGE)
    add_textbox(slide, Inches(4.0), Inches(2.1), Inches(3), Inches(0.4),
                'AWS CLOUD', font_size=14, font_color=ACCENT_ORANGE, bold=True)

    # EC2 box
    ec2_card = add_card(slide, Inches(4.2), Inches(2.8), Inches(3.8), Inches(3.8), MEDIUM_BG, ACCENT_BLUE)
    add_textbox(slide, Inches(4.2), Inches(2.9), Inches(3.8), Inches(0.35),
                'AWS EC2 Instance', font_size=15, font_color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(4.2), Inches(3.25), Inches(3.8), Inches(0.25),
                'Ubuntu Server (t2.micro)', font_size=10, font_color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # Nginx sub-card
    nginx_card = add_card(slide, Inches(4.5), Inches(3.7), Inches(3.2), Inches(0.7), LIGHT_BG, GREEN_ACCENT)
    add_textbox(slide, Inches(4.5), Inches(3.75), Inches(3.2), Inches(0.55),
                'Nginx Reverse Proxy (Port 80)', font_size=12, font_color=GREEN_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

    # Arrow down
    add_shape(slide, MSO_SHAPE.DOWN_ARROW, Inches(5.8), Inches(4.45), Inches(0.4), Inches(0.4), LIGHT_GRAY)

    # Gunicorn sub-card
    gunicorn_card = add_card(slide, Inches(4.5), Inches(4.9), Inches(3.2), Inches(0.7), LIGHT_BG, ACCENT_PURPLE)
    add_textbox(slide, Inches(4.5), Inches(4.95), Inches(3.2), Inches(0.55),
                'Gunicorn WSGI (Port 5000)', font_size=12, font_color=ACCENT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

    # Arrow down
    add_shape(slide, MSO_SHAPE.DOWN_ARROW, Inches(5.8), Inches(5.65), Inches(0.4), Inches(0.4), LIGHT_GRAY)

    # Flask App sub-card
    flask_card = add_card(slide, Inches(4.5), Inches(6.1), Inches(3.2), Inches(0.55), LIGHT_BG, ACCENT_ORANGE)
    add_textbox(slide, Inches(4.5), Inches(6.12), Inches(3.2), Inches(0.45),
                'Flask Application (EmpTrack)', font_size=12, font_color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

    # RDS box
    rds_card = add_card(slide, Inches(8.8), Inches(2.8), Inches(3.5), Inches(2.0), MEDIUM_BG, ACCENT_ORANGE)
    add_textbox(slide, Inches(8.8), Inches(2.9), Inches(3.5), Inches(0.35),
                'AWS RDS', font_size=15, font_color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(8.8), Inches(3.3), Inches(3.5), Inches(1.2),
                'MySQL 8.0 Database\ndb.t3.micro (Free Tier)\n\nStores: Users, Attendance,\nPayroll, Leave, Documents',
                font_size=11, font_color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # S3 box
    s3_card = add_card(slide, Inches(8.8), Inches(5.1), Inches(3.5), Inches(1.7), MEDIUM_BG, ACCENT_TEAL)
    add_textbox(slide, Inches(8.8), Inches(5.2), Inches(3.5), Inches(0.35),
                'AWS S3 Bucket', font_size=15, font_color=ACCENT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(8.8), Inches(5.6), Inches(3.5), Inches(1.0),
                'Document & File Storage\n\nStores: Resumes, PAN,\nAadhaar, Salary Slips',
                font_size=11, font_color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Arrows between EC2 and RDS/S3
    add_shape(slide, MSO_SHAPE.RIGHT_ARROW, Inches(8.1), Inches(3.5), Inches(0.6), Inches(0.3), ACCENT_ORANGE)
    add_shape(slide, MSO_SHAPE.RIGHT_ARROW, Inches(8.1), Inches(5.6), Inches(0.6), Inches(0.3), ACCENT_TEAL)


def slide_tech_stack(prs):
    """Slide 6: Technology Stack"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_decorative_circles(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
                'TECHNOLOGY STACK', font_size=14, font_color=ACCENT_ORANGE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_ORANGE)
    add_textbox(slide, Inches(0.8), Inches(1.15), Inches(10), Inches(0.6),
                'Tools & Technologies Used', font_size=30, font_color=WHITE, bold=True)

    categories = [
        ('Backend', [
            ('Python 3.11', 'Core programming language'),
            ('Flask', 'Lightweight web framework'),
            ('SQLAlchemy', 'ORM for database operations'),
            ('Gunicorn', 'Production WSGI HTTP server'),
        ], ACCENT_BLUE),
        ('Database', [
            ('MySQL 8.0', 'Relational database engine'),
            ('AWS RDS', 'Managed database service'),
            ('PyMySQL', 'Python MySQL connector'),
            ('Flask-Migrate', 'Database migration tool'),
        ], ACCENT_ORANGE),
        ('Frontend', [
            ('HTML5 / Jinja2', 'Templating and structure'),
            ('Vanilla CSS', 'Custom design system'),
            ('JavaScript', 'Dynamic UI interactions'),
            ('Chart.js', 'Analytics & data charts'),
        ], ACCENT_TEAL),
        ('AWS Cloud', [
            ('EC2 (t2.micro)', 'Virtual machine hosting'),
            ('RDS MySQL', 'Managed database service'),
            ('S3 Bucket', 'File/document storage'),
            ('Security Groups', 'Firewall & access control'),
        ], ACCENT_PURPLE),
    ]

    for col_idx, (cat_name, items, color) in enumerate(categories):
        x = Inches(0.5) + col_idx * Inches(3.15)
        card = add_card(slide, x, Inches(2.2), Inches(2.95), Inches(4.8), MEDIUM_BG, color)
        add_textbox(slide, x, Inches(2.35), Inches(2.95), Inches(0.4),
                    cat_name, font_size=18, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_accent_line(slide, x + Inches(0.5), Inches(2.8), Inches(1.95), color, Pt(2))

        for i, (tech, desc) in enumerate(items):
            y = Inches(3.1) + i * Inches(0.9)
            add_textbox(slide, x + Inches(0.25), y, Inches(2.5), Inches(0.3),
                        tech, font_size=14, font_color=WHITE, bold=True)
            add_textbox(slide, x + Inches(0.25), y + Inches(0.3), Inches(2.5), Inches(0.3),
                        desc, font_size=11, font_color=MID_GRAY)


def slide_aws_services_why(prs):
    """Slide 7: Why AWS? Question Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_decorative_circles(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
                'AWS SERVICES', font_size=14, font_color=ACCENT_ORANGE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_ORANGE)
    add_textbox(slide, Inches(0.8), Inches(1.15), Inches(10), Inches(0.6),
                'Why Did We Choose AWS?', font_size=30, font_color=WHITE, bold=True)

    # Question card
    q_card = add_card(slide, Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.2), LIGHT_BG, ACCENT_PURPLE)
    add_textbox(slide, Inches(1.2), Inches(2.15), Inches(10), Inches(0.35),
                'Why not deploy on a regular shared hosting or a local server?', font_size=16, font_color=ACCENT_PURPLE, bold=True)
    add_textbox(slide, Inches(1.2), Inches(2.55), Inches(10.5), Inches(0.5),
                'Shared hosting lacks root access, has limited scalability, and cannot run background processes like Gunicorn. '
                'Local servers require constant uptime, manual maintenance, and are not accessible globally. '
                'AWS provides enterprise-grade infrastructure with pay-as-you-go pricing.',
                font_size=12, font_color=LIGHT_GRAY)

    services = [
        ('Amazon EC2', 'Elastic Compute Cloud',
         'Virtual server (Ubuntu) to host\nour Flask application with Nginx\nand Gunicorn. Provides full root\naccess and SSH connectivity.',
         'Why EC2? We need a full Linux\nserver with root access to install\nPython, Nginx, and run Gunicorn\nas a background service.', ACCENT_BLUE),
        ('Amazon RDS', 'Relational Database Service',
         'Managed MySQL 8.0 database\nwith automated backups, patches,\nand high availability. Stores all\nemployee and payroll data.',
         'Why RDS? Managing MySQL on\nEC2 requires manual backups,\nupdates, and monitoring. RDS\nhandles all of this automatically.', ACCENT_ORANGE),
        ('Amazon S3', 'Simple Storage Service',
         'Scalable object storage for\nemployee documents (Resumes,\nPAN, Aadhaar, Salary PDFs).\n11 nines of durability.',
         'Why S3? Storing files on EC2\nEBS is limited and expensive.\nS3 offers unlimited, durable\nstorage at very low cost.', ACCENT_TEAL),
    ]

    for i, (name, subtitle, desc, why, color) in enumerate(services):
        x = Inches(0.5) + i * Inches(4.2)
        card = add_card(slide, x, Inches(3.6), Inches(3.9), Inches(3.5), MEDIUM_BG, color)
        add_textbox(slide, x, Inches(3.75), Inches(3.9), Inches(0.35),
                    name, font_size=18, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(4.1), Inches(3.9), Inches(0.25),
                    subtitle, font_size=10, font_color=MID_GRAY, alignment=PP_ALIGN.CENTER)
        add_accent_line(slide, x + Inches(0.5), Inches(4.4), Inches(2.9), color, Pt(1.5))
        add_textbox(slide, x + Inches(0.25), Inches(4.6), Inches(3.4), Inches(1.1),
                    desc, font_size=11, font_color=LIGHT_GRAY)
        # Why box
        add_textbox(slide, x + Inches(0.25), Inches(5.75), Inches(3.4), Inches(1.1),
                    why, font_size=10, font_color=MID_GRAY)


def slide_phase1_local(prs):
    """Slide 8: Phase 1 — Local Development"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_decorative_circles(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
                'PHASE 1: LOCAL DEVELOPMENT', font_size=14, font_color=ACCENT_ORANGE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_ORANGE)
    add_textbox(slide, Inches(0.8), Inches(1.15), Inches(10), Inches(0.6),
                'Building the Application Locally', font_size=30, font_color=WHITE, bold=True)

    steps = [
        ('Step 1', 'Setup Project Structure', 'Created Flask app with Blueprint architecture:\nModels, Routes, Services, Templates, Static assets'),
        ('Step 2', 'Database Models', 'Designed 6 SQLAlchemy models:\nUser, Department, Attendance, LeaveRequest,\nLeaveBalance, Payroll, Document, AuditLog'),
        ('Step 3', 'Business Logic', 'Built payroll engine with Indian tax calculations:\nPF (12%), ESI (0.75%), TDS (Income Tax slabs),\nPDF salary slip generator using ReportLab'),
        ('Step 4', 'Frontend & Templates', 'Created 30+ Jinja2 HTML templates with a\ncustom CSS design system, responsive sidebar\nlayout, and Chart.js analytics dashboards'),
        ('Step 5', 'Local Testing', 'Connected to local MySQL database,\nseeded demo data (Admin, HR, Employee accounts),\ntested all routes on http://localhost:5000'),
    ]

    for i, (step, title, desc) in enumerate(steps):
        y = Inches(2.0) + i * Inches(1.05)
        # Step badge
        badge = add_card(slide, Inches(0.8), y, Inches(1.0), Inches(0.5), ACCENT_BLUE)
        add_textbox(slide, Inches(0.8), y + Inches(0.05), Inches(1.0), Inches(0.4),
                    step, font_size=11, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        # Title
        add_textbox(slide, Inches(2.0), y, Inches(3), Inches(0.35),
                    title, font_size=16, font_color=WHITE, bold=True)
        # Description
        add_textbox(slide, Inches(2.0), y + Inches(0.35), Inches(10), Inches(0.7),
                    desc, font_size=12, font_color=LIGHT_GRAY)


def slide_phase2_ec2(prs):
    """Slide 9: Phase 2 — EC2 Setup"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_decorative_circles(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
                'PHASE 2: AWS EC2 SETUP', font_size=14, font_color=ACCENT_ORANGE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_ORANGE)
    add_textbox(slide, Inches(0.8), Inches(1.15), Inches(10), Inches(0.6),
                'Launching the Cloud Server', font_size=30, font_color=WHITE, bold=True)

    # Question
    q_card = add_card(slide, Inches(0.8), Inches(2.0), Inches(11.7), Inches(0.8), LIGHT_BG, ACCENT_PURPLE)
    add_textbox(slide, Inches(1.2), Inches(2.1), Inches(10), Inches(0.3),
                'What is Amazon EC2?', font_size=15, font_color=ACCENT_PURPLE, bold=True)
    add_textbox(slide, Inches(1.2), Inches(2.4), Inches(10.5), Inches(0.3),
                'Elastic Compute Cloud — A virtual server in the cloud. It gives you full control over the operating system, software, and network configuration.',
                font_size=12, font_color=LIGHT_GRAY)

    steps = [
        ('1', 'Go to AWS Console > EC2 > Launch Instance'),
        ('2', 'Name the instance: EmpTrack-Server'),
        ('3', 'Select AMI: Ubuntu Server 22.04 LTS (64-bit)'),
        ('4', 'Instance Type: t2.micro (Free Tier Eligible)'),
        ('5', 'Create Key Pair: emptrack-key.pem (RSA)'),
        ('6', 'Security Group: Allow SSH (22), HTTP (80), Custom TCP (5000)'),
        ('7', 'Storage: 8 GB General Purpose SSD (gp3)'),
        ('8', 'Click Launch Instance > Copy Public IPv4 Address'),
    ]

    for i, (num, text) in enumerate(steps):
        col = i // 4
        row = i % 4
        x = Inches(0.8) + col * Inches(6.2)
        y = Inches(3.2) + row * Inches(0.95)
        badge = add_shape(slide, MSO_SHAPE.OVAL, x, y, Inches(0.4), Inches(0.4), ACCENT_BLUE)
        add_textbox(slide, x, y + Inches(0.02), Inches(0.4), Inches(0.35),
                    num, font_size=13, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.55), y + Inches(0.05), Inches(5.5), Inches(0.35),
                    text, font_size=14, font_color=LIGHT_GRAY)

    # SSH Command
    cmd_card = add_card(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.6), RGBColor(0x0A, 0x14, 0x22))
    add_textbox(slide, Inches(1.2), Inches(6.55), Inches(10), Inches(0.45),
                'ssh -i "emptrack-key.pem" ubuntu@<EC2_PUBLIC_IP>',
                font_size=14, font_color=GREEN_ACCENT, font_name='Consolas')


def slide_phase3_rds(prs):
    """Slide 10: Phase 3 — RDS Setup"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_decorative_circles(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
                'PHASE 3: AWS RDS SETUP', font_size=14, font_color=ACCENT_ORANGE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_ORANGE)
    add_textbox(slide, Inches(0.8), Inches(1.15), Inches(10), Inches(0.6),
                'Creating the Managed MySQL Database', font_size=30, font_color=WHITE, bold=True)

    # Question
    q_card = add_card(slide, Inches(0.8), Inches(2.0), Inches(11.7), Inches(0.8), LIGHT_BG, ACCENT_PURPLE)
    add_textbox(slide, Inches(1.2), Inches(2.1), Inches(10), Inches(0.3),
                'Why RDS instead of installing MySQL directly on EC2?', font_size=15, font_color=ACCENT_PURPLE, bold=True)
    add_textbox(slide, Inches(1.2), Inches(2.4), Inches(10.5), Inches(0.3),
                'RDS provides automated backups, automatic patching, multi-AZ failover, and monitoring — all managed by AWS. No manual DB administration needed.',
                font_size=12, font_color=LIGHT_GRAY)

    steps = [
        ('1', 'Go to AWS Console > RDS > Create Database'),
        ('2', 'Engine: MySQL 8.0 (Community Edition)'),
        ('3', 'Template: Free Tier'),
        ('4', 'DB Instance: emptrack-db-instance'),
        ('5', 'Master User: admin / Strong Password'),
        ('6', 'Instance: db.t3.micro (Free Tier)'),
        ('7', 'Initial Database Name: employee_db'),
        ('8', 'Security Group: emptrack-rds-sg'),
    ]

    for i, (num, text) in enumerate(steps):
        col = i // 4
        row = i % 4
        x = Inches(0.8) + col * Inches(6.2)
        y = Inches(3.2) + row * Inches(0.85)
        badge = add_shape(slide, MSO_SHAPE.OVAL, x, y, Inches(0.4), Inches(0.4), ACCENT_ORANGE)
        add_textbox(slide, x, y + Inches(0.02), Inches(0.4), Inches(0.35),
                    num, font_size=13, font_color=DARK_BG, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.55), y + Inches(0.05), Inches(5.5), Inches(0.35),
                    text, font_size=14, font_color=LIGHT_GRAY)

    # Connection string
    cmd_card = add_card(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.6), RGBColor(0x0A, 0x14, 0x22))
    add_textbox(slide, Inches(1.2), Inches(5.85), Inches(10), Inches(0.45),
                'DATABASE_URL=mysql+pymysql://admin:password@emptrack-db-instance.xxxxx.rds.amazonaws.com/employee_db',
                font_size=12, font_color=GREEN_ACCENT, font_name='Consolas')

    # Security Group Config
    sg_card = add_card(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.7), MEDIUM_BG, ACCENT_ORANGE)
    add_textbox(slide, Inches(1.2), Inches(6.55), Inches(4), Inches(0.3),
                'Security Group Rule:', font_size=13, font_color=ACCENT_ORANGE, bold=True)
    add_textbox(slide, Inches(1.2), Inches(6.85), Inches(10), Inches(0.3),
                'Inbound Rule: MYSQL/Aurora (Port 3306) > Source: emptrack-ec2-sg (allows EC2 to connect to RDS)',
                font_size=12, font_color=LIGHT_GRAY)


def slide_phase4_s3(prs):
    """Slide 11: Phase 4 — S3 Setup"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_decorative_circles(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
                'PHASE 4: AWS S3 SETUP', font_size=14, font_color=ACCENT_ORANGE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_ORANGE)
    add_textbox(slide, Inches(0.8), Inches(1.15), Inches(10), Inches(0.6),
                'Setting Up Document & File Storage', font_size=30, font_color=WHITE, bold=True)

    # Question
    q_card = add_card(slide, Inches(0.8), Inches(2.0), Inches(11.7), Inches(0.8), LIGHT_BG, ACCENT_PURPLE)
    add_textbox(slide, Inches(1.2), Inches(2.1), Inches(10), Inches(0.3),
                'What is Amazon S3 and why do we need it?', font_size=15, font_color=ACCENT_PURPLE, bold=True)
    add_textbox(slide, Inches(1.2), Inches(2.4), Inches(10.5), Inches(0.3),
                'S3 (Simple Storage Service) provides unlimited, highly durable object storage. Employee documents like Resumes, PAN cards, and PDF salary slips are stored here instead of on the EC2 instance.',
                font_size=12, font_color=LIGHT_GRAY)

    # Steps
    steps = [
        ('1', 'Go to AWS Console > S3 > Create Bucket'),
        ('2', 'Bucket Name: emptrack-storage-bucket-2024'),
        ('3', 'Region: ap-south-1 (Mumbai) or us-east-1'),
        ('4', 'Block Public Access: Keep enabled for security'),
        ('5', 'Versioning: Disabled (or enable for audit trail)'),
        ('6', 'Click Create Bucket'),
    ]

    for i, (num, text) in enumerate(steps):
        col = i // 3
        row = i % 3
        x = Inches(0.8) + col * Inches(6.2)
        y = Inches(3.2) + row * Inches(0.85)
        badge = add_shape(slide, MSO_SHAPE.OVAL, x, y, Inches(0.4), Inches(0.4), ACCENT_TEAL)
        add_textbox(slide, x, y + Inches(0.02), Inches(0.4), Inches(0.35),
                    num, font_size=13, font_color=DARK_BG, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.55), y + Inches(0.05), Inches(5.5), Inches(0.35),
                    text, font_size=14, font_color=LIGHT_GRAY)

    # What gets stored
    store_card = add_card(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.6), MEDIUM_BG, ACCENT_TEAL)
    add_textbox(slide, Inches(1.2), Inches(5.6), Inches(4), Inches(0.35),
                'What Gets Stored in S3?', font_size=16, font_color=ACCENT_TEAL, bold=True)

    doc_types = [
        ('Employee Resumes', 'PDF/DOC uploads'),
        ('Aadhaar & PAN Cards', 'Identity documents'),
        ('Offer Letters', 'Employment documents'),
        ('Salary Slips (PDF)', 'Auto-generated payslips'),
        ('Profile Photos', 'Employee avatars'),
        ('Certificates', 'Training & education'),
    ]
    for i, (doc, desc) in enumerate(doc_types):
        col = i // 2
        row = i % 2
        x = Inches(1.2) + col * Inches(3.6)
        y = Inches(6.05) + row * Inches(0.45)
        add_textbox(slide, x, y, Inches(1.8), Inches(0.3),
                    doc, font_size=12, font_color=WHITE, bold=True)
        add_textbox(slide, x + Inches(1.8), y, Inches(1.5), Inches(0.3),
                    desc, font_size=11, font_color=MID_GRAY)


def slide_phase5_deployment(prs):
    """Slide 12: Phase 5 — Deployment"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_decorative_circles(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
                'PHASE 5: DEPLOYMENT & GO-LIVE', font_size=14, font_color=ACCENT_ORANGE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_ORANGE)
    add_textbox(slide, Inches(0.8), Inches(1.15), Inches(10), Inches(0.6),
                'Deploying on EC2 with Nginx & Gunicorn', font_size=30, font_color=WHITE, bold=True)

    commands = [
        ('1. Upload Code to EC2', 'scp -i "emptrack-key.pem" employee_system.zip ubuntu@<IP>:~/'),
        ('2. Install Dependencies', 'sudo apt update && sudo apt install python3-pip python3-venv nginx -y'),
        ('3. Setup Virtual Env', 'python3 -m venv venv && source venv/bin/activate\npip install -r requirements.txt && pip install gunicorn'),
        ('4. Configure .env', 'nano .env\n# Set DATABASE_URL to RDS endpoint\n# Set SECRET_KEY for production'),
        ('5. Seed Database', 'python3 seed.py\n# Creates tables on RDS & inserts demo data'),
        ('6. Create Systemd Service', 'sudo nano /etc/systemd/system/emptrack.service\nsudo systemctl start emptrack && sudo systemctl enable emptrack'),
        ('7. Configure Nginx', 'sudo nano /etc/nginx/sites-available/emptrack\n# Reverse proxy port 80 -> 127.0.0.1:5000\nsudo systemctl restart nginx'),
        ('8. Go Live!', 'Visit http://<EC2_PUBLIC_IP>\n# Application is now accessible worldwide!'),
    ]

    for i, (title, cmd) in enumerate(commands):
        col = i // 4
        row = i % 4
        x = Inches(0.5) + col * Inches(6.3)
        y = Inches(2.0) + row * Inches(1.3)
        # Title
        add_textbox(slide, x, y, Inches(5.8), Inches(0.3),
                    title, font_size=14, font_color=ACCENT_ORANGE, bold=True)
        # Command block
        cmd_card = add_card(slide, x, y + Inches(0.35), Inches(5.8), Inches(0.85), RGBColor(0x0A, 0x14, 0x22))
        add_textbox(slide, x + Inches(0.15), y + Inches(0.38), Inches(5.5), Inches(0.75),
                    cmd, font_size=10, font_color=GREEN_ACCENT, font_name='Consolas')


def slide_security(prs):
    """Slide 13: Security & Best Practices"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_decorative_circles(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
                'SECURITY & BEST PRACTICES', font_size=14, font_color=ACCENT_ORANGE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_ORANGE)
    add_textbox(slide, Inches(0.8), Inches(1.15), Inches(10), Inches(0.6),
                'How We Secured the Application', font_size=30, font_color=WHITE, bold=True)

    security_items = [
        ('AWS Security Groups', 'Configured firewall rules to allow only\nnecessary traffic (SSH, HTTP, MySQL)\nbetween EC2, RDS, and the internet.', ACCENT_BLUE),
        ('VPC Network Isolation', 'RDS is accessible only from EC2 via\nprivate networking within the same VPC.\nNo direct public internet access to DB.', ACCENT_ORANGE),
        ('Password Hashing', 'All user passwords are hashed using\nWerkzeug\'s generate_password_hash()\nbefore storing in the database.', ACCENT_TEAL),
        ('Role-Based Access Control', 'Flask-Login enforces authentication.\nAdmin, HR, and Employee roles have\nspecific permissions and route guards.', ACCENT_PURPLE),
        ('Environment Variables', '.env file stores secrets (DB password,\nSECRET_KEY). Never committed to Git.\nIncluded in .gitignore.', GREEN_ACCENT),
        ('SSH Key Pair Auth', 'EC2 access uses RSA key pair (.pem file)\ninstead of password authentication.\nKey pair stored securely on local machine.', RED_ACCENT),
    ]

    for i, (title, desc, color) in enumerate(security_items):
        col = i % 3
        row = i // 3
        x = Inches(0.5) + col * Inches(4.2)
        y = Inches(2.2) + row * Inches(2.5)
        card = add_card(slide, x, y, Inches(3.9), Inches(2.2), MEDIUM_BG, color)
        add_textbox(slide, x, y + Inches(0.15), Inches(3.9), Inches(0.35),
                    title, font_size=16, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_accent_line(slide, x + Inches(0.5), y + Inches(0.55), Inches(2.9), color, Pt(1.5))
        add_textbox(slide, x + Inches(0.3), y + Inches(0.75), Inches(3.3), Inches(1.2),
                    desc, font_size=12, font_color=LIGHT_GRAY)


def slide_security_groups_detail(prs):
    """Slide 14: Security Groups Detail"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_decorative_circles(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
                'SECURITY GROUPS CONFIGURATION', font_size=14, font_color=ACCENT_ORANGE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_ORANGE)
    add_textbox(slide, Inches(0.8), Inches(1.15), Inches(10), Inches(0.6),
                'Firewall Rules for EC2 & RDS', font_size=30, font_color=WHITE, bold=True)

    # Question
    q_card = add_card(slide, Inches(0.8), Inches(2.0), Inches(11.7), Inches(0.8), LIGHT_BG, ACCENT_PURPLE)
    add_textbox(slide, Inches(1.2), Inches(2.1), Inches(10), Inches(0.3),
                'What are Security Groups in AWS?', font_size=15, font_color=ACCENT_PURPLE, bold=True)
    add_textbox(slide, Inches(1.2), Inches(2.4), Inches(10.5), Inches(0.3),
                'Security Groups act as virtual firewalls for your AWS resources. They control inbound and outbound traffic at the instance level using allow rules.',
                font_size=12, font_color=LIGHT_GRAY)

    # EC2 Security Group Table
    add_textbox(slide, Inches(0.8), Inches(3.1), Inches(5), Inches(0.35),
                'EC2 Security Group: emptrack-ec2-sg', font_size=15, font_color=ACCENT_BLUE, bold=True)

    ec2_rules = [
        ('Type', 'Protocol', 'Port', 'Source', 'Purpose'),
        ('SSH', 'TCP', '22', 'My IP', 'Remote Access'),
        ('HTTP', 'TCP', '80', '0.0.0.0/0', 'Web Access'),
        ('Custom TCP', 'TCP', '5000', '0.0.0.0/0', 'Flask Port'),
    ]

    for i, row in enumerate(ec2_rules):
        y = Inches(3.5) + i * Inches(0.4)
        bg_color = LIGHT_BG if i == 0 else (MEDIUM_BG if i % 2 == 0 else DARK_BG)
        text_color = ACCENT_BLUE if i == 0 else LIGHT_GRAY
        font_bold = i == 0
        for j, cell in enumerate(row):
            x = Inches(0.8) + j * Inches(2.0)
            w = Inches(1.9)
            cell_card = add_card(slide, x, y, w, Inches(0.38), bg_color)
            add_textbox(slide, x + Inches(0.1), y + Inches(0.02), w - Inches(0.2), Inches(0.32),
                        cell, font_size=11, font_color=text_color, bold=font_bold, alignment=PP_ALIGN.CENTER)

    # RDS Security Group Table
    add_textbox(slide, Inches(0.8), Inches(5.4), Inches(5), Inches(0.35),
                'RDS Security Group: emptrack-rds-sg', font_size=15, font_color=ACCENT_ORANGE, bold=True)

    rds_rules = [
        ('Type', 'Protocol', 'Port', 'Source', 'Purpose'),
        ('MYSQL/Aurora', 'TCP', '3306', 'emptrack-ec2-sg', 'EC2 -> RDS'),
    ]

    for i, row in enumerate(rds_rules):
        y = Inches(5.8) + i * Inches(0.4)
        bg_color = LIGHT_BG if i == 0 else MEDIUM_BG
        text_color = ACCENT_ORANGE if i == 0 else LIGHT_GRAY
        font_bold = i == 0
        for j, cell in enumerate(row):
            x = Inches(0.8) + j * Inches(2.0)
            w = Inches(1.9)
            cell_card = add_card(slide, x, y, w, Inches(0.38), bg_color)
            add_textbox(slide, x + Inches(0.1), y + Inches(0.02), w - Inches(0.2), Inches(0.32),
                        cell, font_size=11, font_color=text_color, bold=font_bold, alignment=PP_ALIGN.CENTER)


def slide_app_features(prs):
    """Slide 15: Application Features"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_decorative_circles(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
                'APPLICATION FEATURES', font_size=14, font_color=ACCENT_ORANGE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_ORANGE)
    add_textbox(slide, Inches(0.8), Inches(1.15), Inches(10), Inches(0.6),
                'Key Modules of EmpTrack', font_size=30, font_color=WHITE, bold=True)

    modules = [
        ('QR Code Attendance', 'Admin generates daily QR code.\nEmployees scan to check in/out.\nPrevents proxy attendance.', ACCENT_BLUE),
        ('GPS Verification', 'Browser Geolocation API captures\nemployee coordinates.\nAttendance accepted only if in range.', ACCENT_ORANGE),
        ('Automated Payroll', 'Indian tax slabs, PF (12%),\nESI (0.75%), TDS calculation.\nAuto-generated monthly payroll.', ACCENT_TEAL),
        ('PDF Salary Slips', 'ReportLab generates professional\nprint-ready salary slips with\ncomplete earnings & deductions.', ACCENT_PURPLE),
        ('Leave Management', 'Casual, Sick, Earned leave tracking.\nHR approval workflow.\nAuto-deduction from payroll.', GREEN_ACCENT),
        ('Document Storage', 'Upload Resumes, PAN, Aadhaar.\nStored securely on AWS S3.\nAdmin/HR download access.', ACCENT_BLUE),
        ('Role-Based Dashboards', 'Admin: Full control & analytics.\nHR: Attendance & leave mgmt.\nEmployee: Self-service portal.', ACCENT_ORANGE),
        ('Audit Logging', 'Every login, action, and change\nis recorded in audit trail.\nAdmin can review all system activity.', RED_ACCENT),
    ]

    for i, (title, desc, color) in enumerate(modules):
        col = i % 4
        row = i // 4
        x = Inches(0.4) + col * Inches(3.2)
        y = Inches(2.2) + row * Inches(2.5)
        card = add_card(slide, x, y, Inches(3.0), Inches(2.2), MEDIUM_BG, color)
        add_textbox(slide, x, y + Inches(0.15), Inches(3.0), Inches(0.35),
                    title, font_size=14, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_accent_line(slide, x + Inches(0.4), y + Inches(0.55), Inches(2.2), color, Pt(1.5))
        add_textbox(slide, x + Inches(0.25), y + Inches(0.75), Inches(2.5), Inches(1.2),
                    desc, font_size=11, font_color=LIGHT_GRAY)


def slide_database_schema(prs):
    """Slide 16: Database Schema"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_decorative_circles(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
                'DATABASE DESIGN', font_size=14, font_color=ACCENT_ORANGE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_ORANGE)
    add_textbox(slide, Inches(0.8), Inches(1.15), Inches(10), Inches(0.6),
                'MySQL Database Schema on AWS RDS', font_size=30, font_color=WHITE, bold=True)

    tables = [
        ('users', 'id, employee_id, name, email,\npassword_hash, role, department_id,\ndesignation, base_salary, phone,\njoin_date, pan_number, bank_account', ACCENT_BLUE),
        ('departments', 'id, name, description,\ncreated_at', ACCENT_ORANGE),
        ('attendance', 'id, employee_id, date, check_in,\ncheck_out, status, method,\nwork_hours, overtime_hours', ACCENT_TEAL),
        ('leave_requests', 'id, employee_id, leave_type,\nfrom_date, to_date, days,\nreason, status, approved_by', ACCENT_PURPLE),
        ('leave_balances', 'id, employee_id, casual_leave,\nsick_leave, earned_leave', GREEN_ACCENT),
        ('payroll', 'id, employee_id, month, year,\nbase_salary, gross_earnings,\npf, esi, tds, net_salary', ACCENT_ORANGE),
        ('documents', 'id, employee_id, doc_type,\noriginal_filename, stored_filename,\nfile_path, file_size, mime_type', ACCENT_BLUE),
        ('audit_logs', 'id, user_id, action, target_type,\ntarget_id, details, ip_address,\ntimestamp', RED_ACCENT),
    ]

    for i, (table, columns, color) in enumerate(tables):
        col = i % 4
        row = i // 4
        x = Inches(0.4) + col * Inches(3.2)
        y = Inches(2.2) + row * Inches(2.5)
        card = add_card(slide, x, y, Inches(3.0), Inches(2.2), MEDIUM_BG, color)
        add_textbox(slide, x, y + Inches(0.1), Inches(3.0), Inches(0.35),
                    table, font_size=14, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_accent_line(slide, x + Inches(0.4), y + Inches(0.5), Inches(2.2), color, Pt(1.5))
        add_textbox(slide, x + Inches(0.2), y + Inches(0.7), Inches(2.6), Inches(1.3),
                    columns, font_size=10, font_color=LIGHT_GRAY, font_name='Consolas')


def slide_project_timeline(prs):
    """Slide 17: Project Timeline"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_decorative_circles(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
                'PROJECT TIMELINE', font_size=14, font_color=ACCENT_ORANGE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_ORANGE)
    add_textbox(slide, Inches(0.8), Inches(1.15), Inches(10), Inches(0.6),
                'Development Phases & Milestones', font_size=30, font_color=WHITE, bold=True)

    phases = [
        ('Week 1-2', 'Planning & Research', 'Requirements gathering, architecture\ndesign, AWS service selection,\ntechnology stack finalization', ACCENT_BLUE),
        ('Week 3-4', 'Backend Development', 'Flask app factory, SQLAlchemy models,\npayroll engine with Indian tax logic,\nQR service, PDF generator', ACCENT_ORANGE),
        ('Week 5-6', 'Frontend & Templates', '30+ Jinja2 templates, custom CSS\ndesign system, responsive layouts,\nChart.js dashboards, JS interactions', ACCENT_TEAL),
        ('Week 7', 'AWS Infrastructure', 'EC2 launch, RDS MySQL setup,\nS3 bucket creation, security groups\nconfiguration, VPC networking', ACCENT_PURPLE),
        ('Week 8', 'Deployment & Testing', 'Nginx + Gunicorn setup, systemd\nservice, database seeding on RDS,\nend-to-end testing, go-live', GREEN_ACCENT),
    ]

    # Timeline line
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(3.85), Inches(10.5), Pt(3), ACCENT_ORANGE)

    for i, (time, title, desc, color) in enumerate(phases):
        x = Inches(0.8) + i * Inches(2.45)
        # Circle on timeline
        dot = add_shape(slide, MSO_SHAPE.OVAL, x + Inches(0.8), Inches(3.7), Inches(0.35), Inches(0.35), color)
        # Time label above
        add_textbox(slide, x, Inches(2.8), Inches(2.3), Inches(0.35),
                    time, font_size=14, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(3.15), Inches(2.3), Inches(0.35),
                    title, font_size=12, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        # Description below
        card = add_card(slide, x, Inches(4.4), Inches(2.3), Inches(2.2), MEDIUM_BG, color)
        add_textbox(slide, x + Inches(0.15), Inches(4.55), Inches(2.0), Inches(1.8),
                    desc, font_size=11, font_color=LIGHT_GRAY)


def slide_demo(prs):
    """Slide 18: Live Demo"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_decorative_circles(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
                'LIVE DEMO', font_size=14, font_color=ACCENT_ORANGE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_ORANGE)
    add_textbox(slide, Inches(0.8), Inches(1.15), Inches(10), Inches(0.6),
                'Working Application Walkthrough', font_size=30, font_color=WHITE, bold=True)

    add_textbox(slide, Inches(1.5), Inches(2.3), Inches(10), Inches(0.5),
                'Access URL:  http://<EC2_PUBLIC_IP>',
                font_size=22, font_color=GREEN_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

    # Demo credentials
    creds = [
        ('Admin', 'admin@emptrack.com', 'Admin@123', 'Full system control, analytics, payroll processing', ACCENT_BLUE),
        ('HR', 'hr@emptrack.com', 'Hr@123', 'Attendance management, leave approvals', ACCENT_ORANGE),
        ('Employee', 'emp@emptrack.com', 'Emp@123', 'Self-service: attendance, leave, salary slips', ACCENT_TEAL),
    ]

    for i, (role, email, pwd, desc, color) in enumerate(creds):
        x = Inches(0.5) + i * Inches(4.2)
        card = add_card(slide, x, Inches(3.2), Inches(3.9), Inches(2.5), MEDIUM_BG, color)
        add_textbox(slide, x, Inches(3.35), Inches(3.9), Inches(0.35),
                    f'{role} Login', font_size=18, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_accent_line(slide, x + Inches(0.5), Inches(3.75), Inches(2.9), color, Pt(1.5))
        add_textbox(slide, x + Inches(0.3), Inches(3.95), Inches(3.3), Inches(0.3),
                    f'Email: {email}', font_size=13, font_color=WHITE, font_name='Consolas')
        add_textbox(slide, x + Inches(0.3), Inches(4.25), Inches(3.3), Inches(0.3),
                    f'Password: {pwd}', font_size=13, font_color=WHITE, font_name='Consolas')
        add_textbox(slide, x + Inches(0.3), Inches(4.7), Inches(3.3), Inches(0.6),
                    desc, font_size=11, font_color=MID_GRAY)

    # Demo walkthrough
    demo_steps = [
        'Login as Admin > View Dashboard with analytics charts',
        'Navigate to Attendance > Generate QR Code for office',
        'Login as Employee > Scan QR to mark attendance',
        'Apply for Leave > Login as HR > Approve/Reject',
        'Admin: Run Payroll > Employee: Download PDF Salary Slip',
        'Upload Employee Documents > View in Documents section',
    ]

    add_textbox(slide, Inches(0.8), Inches(5.9), Inches(3), Inches(0.35),
                'Demo Walkthrough:', font_size=14, font_color=ACCENT_ORANGE, bold=True)
    for i, step in enumerate(demo_steps):
        col = i // 3
        row = i % 3
        x = Inches(0.8) + col * Inches(6.2)
        y = Inches(6.3) + row * Inches(0.35)
        add_textbox(slide, x, y, Inches(6), Inches(0.3),
                    f'{i+1}. {step}', font_size=11, font_color=LIGHT_GRAY)


def slide_challenges(prs):
    """Slide 19: Challenges & Solutions"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_decorative_circles(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
                'CHALLENGES & SOLUTIONS', font_size=14, font_color=ACCENT_ORANGE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_ORANGE)
    add_textbox(slide, Inches(0.8), Inches(1.15), Inches(10), Inches(0.6),
                'Problems We Faced & How We Solved Them', font_size=30, font_color=WHITE, bold=True)

    challenges = [
        ('MySQL Auth Error', 'MySQL 8+ uses caching_sha2_password\nauthentication by default.', 'Installed the "cryptography" Python\npackage to support SHA256 auth.', RED_ACCENT, GREEN_ACCENT),
        ('Ambiguous Foreign Keys', 'SQLAlchemy threw AmbiguousForeignKeysError\non the Document model (2 FK to users).', 'Added explicit foreign_keys parameter\nto the db.relationship() definition.', RED_ACCENT, GREEN_ACCENT),
        ('Security Group Config', 'EC2 could not connect to RDS.\nConnection timed out on port 3306.', 'Added inbound rule on RDS security group\nto allow traffic from EC2 security group.', RED_ACCENT, GREEN_ACCENT),
        ('Special Chars in Password', 'DB password with "@" broke the\nconnection string URL format.', 'URL-encoded special characters.\n@ becomes %40 in the connection string.', RED_ACCENT, GREEN_ACCENT),
    ]

    for i, (title, problem, solution, p_color, s_color) in enumerate(challenges):
        y = Inches(2.2) + i * Inches(1.25)
        add_textbox(slide, Inches(0.8), y, Inches(3), Inches(0.3),
                    title, font_size=15, font_color=WHITE, bold=True)
        # Problem
        p_card = add_card(slide, Inches(0.8), y + Inches(0.35), Inches(5.5), Inches(0.75), MEDIUM_BG, p_color)
        add_textbox(slide, Inches(1.0), y + Inches(0.38), Inches(1.2), Inches(0.25),
                    'Problem:', font_size=11, font_color=p_color, bold=True)
        add_textbox(slide, Inches(2.2), y + Inches(0.38), Inches(3.8), Inches(0.65),
                    problem, font_size=11, font_color=LIGHT_GRAY)
        # Solution
        s_card = add_card(slide, Inches(6.5), y + Inches(0.35), Inches(6.0), Inches(0.75), MEDIUM_BG, s_color)
        add_textbox(slide, Inches(6.7), y + Inches(0.38), Inches(1.2), Inches(0.25),
                    'Solution:', font_size=11, font_color=s_color, bold=True)
        add_textbox(slide, Inches(7.9), y + Inches(0.38), Inches(4.3), Inches(0.65),
                    solution, font_size=11, font_color=LIGHT_GRAY)


def slide_future_scope(prs):
    """Slide 20: Future Scope"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_decorative_circles(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
                'FUTURE SCOPE', font_size=14, font_color=ACCENT_ORANGE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_ORANGE)
    add_textbox(slide, Inches(0.8), Inches(1.15), Inches(10), Inches(0.6),
                'What Can Be Added Next', font_size=30, font_color=WHITE, bold=True)

    items = [
        ('Face Recognition Attendance', 'Use OpenCV + AWS Rekognition for\nbiometric face-based attendance.\nEliminates all proxy attendance.', ACCENT_BLUE),
        ('Auto-Scaling with ELB', 'Add Elastic Load Balancer and\nAuto Scaling Group for handling\ntraffic spikes during peak hours.', ACCENT_ORANGE),
        ('Email/SMS Notifications', 'AWS SES for email and SNS for SMS\nnotifications on leave approvals,\npayroll generation, and alerts.', ACCENT_TEAL),
        ('CI/CD Pipeline', 'AWS CodePipeline + CodeDeploy for\nautomated deployments from GitHub.\nZero-downtime rolling updates.', ACCENT_PURPLE),
        ('CloudWatch Monitoring', 'Set up AWS CloudWatch for\nperformance monitoring, log analysis,\nand custom metric alarms.', GREEN_ACCENT),
        ('Mobile Application', 'React Native or Flutter mobile app\nfor attendance marking on-the-go\nwith push notifications.', ACCENT_BLUE),
    ]

    for i, (title, desc, color) in enumerate(items):
        col = i % 3
        row = i // 3
        x = Inches(0.5) + col * Inches(4.2)
        y = Inches(2.2) + row * Inches(2.5)
        card = add_card(slide, x, y, Inches(3.9), Inches(2.2), MEDIUM_BG, color)
        add_textbox(slide, x, y + Inches(0.15), Inches(3.9), Inches(0.35),
                    title, font_size=15, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_accent_line(slide, x + Inches(0.5), y + Inches(0.55), Inches(2.9), color, Pt(1.5))
        add_textbox(slide, x + Inches(0.3), y + Inches(0.75), Inches(3.3), Inches(1.2),
                    desc, font_size=12, font_color=LIGHT_GRAY)


def slide_conclusion(prs):
    """Slide 21: Conclusion"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_decorative_circles(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
                'CONCLUSION', font_size=14, font_color=ACCENT_ORANGE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_ORANGE)
    add_textbox(slide, Inches(0.8), Inches(1.15), Inches(10), Inches(0.6),
                'Summary & Key Takeaways', font_size=30, font_color=WHITE, bold=True)

    takeaways = [
        ('Successfully Deployed on AWS', 'The EmpTrack application is live on AWS EC2 with RDS MySQL backend and S3 storage, accessible worldwide via public IP.'),
        ('End-to-End Automation', 'From QR-based attendance marking to automated payroll with Indian tax calculations and PDF salary slip generation — fully automated.'),
        ('Production-Grade Architecture', 'Nginx reverse proxy, Gunicorn WSGI server, systemd service, VPC security groups — enterprise-level deployment stack.'),
        ('Scalable & Secure', 'The architecture supports horizontal scaling with AWS Auto Scaling Groups, Load Balancers, and multi-AZ RDS failover.'),
        ('Hands-On Cloud Experience', 'This project provided practical experience with AWS EC2, RDS, S3, IAM, Security Groups, VPC networking, and Linux server administration.'),
    ]

    for i, (title, desc) in enumerate(takeaways):
        y = Inches(2.1) + i * Inches(1.0)
        badge = add_shape(slide, MSO_SHAPE.OVAL, Inches(0.8), y, Inches(0.4), Inches(0.4), GREEN_ACCENT)
        add_textbox(slide, Inches(0.8), y + Inches(0.02), Inches(0.4), Inches(0.35),
                    str(i+1), font_size=13, font_color=DARK_BG, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.4), y, Inches(4), Inches(0.3),
                    title, font_size=15, font_color=WHITE, bold=True)
        add_textbox(slide, Inches(1.4), y + Inches(0.32), Inches(11), Inches(0.55),
                    desc, font_size=12, font_color=LIGHT_GRAY)


def slide_thankyou(prs):
    """Slide 22: Thank You"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_decorative_circles(slide)

    # Big thank you
    add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
                'Thank You!', font_size=54, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_accent_line(slide, Inches(5), Inches(3.0), Inches(3.3), ACCENT_ORANGE)

    add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.6),
                'Cloud-Based Smart Employee Attendance & Payroll Management System',
                font_size=20, font_color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Project details card
    card = add_card(slide, Inches(3.5), Inches(4.5), Inches(6.3), Inches(2.0), MEDIUM_BG, ACCENT_ORANGE)
    details = [
        ('Project', 'EmpTrack — Employee Attendance & Payroll System'),
        ('Deployed On', 'AWS EC2 + RDS MySQL + S3 Bucket'),
        ('Tech Stack', 'Python Flask, MySQL, Nginx, Gunicorn'),
        ('Presented By', 'Chetan'),
        ('Internship', 'Cloud Computing — August 2026'),
    ]
    for i, (label, value) in enumerate(details):
        y = Inches(4.6) + i * Inches(0.35)
        add_textbox(slide, Inches(3.8), y, Inches(2), Inches(0.3),
                    label + ':', font_size=12, font_color=ACCENT_ORANGE, bold=True)
        add_textbox(slide, Inches(5.8), y, Inches(3.8), Inches(0.3),
                    value, font_size=12, font_color=LIGHT_GRAY)

    add_textbox(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.4),
                'Questions?', font_size=24, font_color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════
# MAIN
# ═════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Build all slides
    slide_title(prs)              # 1. Title
    slide_agenda(prs)             # 2. Agenda
    slide_project_overview(prs)   # 3. Project Overview
    slide_objectives(prs)         # 4. Objectives
    slide_architecture(prs)       # 5. Architecture Diagram
    slide_tech_stack(prs)         # 6. Tech Stack
    slide_aws_services_why(prs)   # 7. AWS Services & Why
    slide_phase1_local(prs)       # 8. Phase 1: Local Dev
    slide_phase2_ec2(prs)         # 9. Phase 2: EC2 Setup
    slide_phase3_rds(prs)         # 10. Phase 3: RDS Setup
    slide_phase4_s3(prs)          # 11. Phase 4: S3 Setup
    slide_phase5_deployment(prs)  # 12. Phase 5: Deployment
    slide_security(prs)           # 13. Security Practices
    slide_security_groups_detail(prs)  # 14. Security Groups Detail
    slide_app_features(prs)       # 15. App Features
    slide_database_schema(prs)    # 16. Database Schema
    slide_project_timeline(prs)   # 17. Timeline
    slide_demo(prs)               # 18. Live Demo
    slide_challenges(prs)         # 19. Challenges
    slide_future_scope(prs)       # 20. Future Scope
    slide_conclusion(prs)         # 21. Conclusion
    slide_thankyou(prs)           # 22. Thank You

    output_path = os.path.join(os.path.dirname(__file__), 'EmpTrack_Cloud_Computing_Internship.pptx')
    prs.save(output_path)
    print(f'Presentation saved to: {output_path}')
    print(f'Total slides: {len(prs.slides)}')


if __name__ == '__main__':
    main()
